from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path(__file__).resolve().parent.parent / "Bazel-Techtalk-30min.pptx"

NAVY = RGBColor(15, 23, 42)
SLATE = RGBColor(51, 65, 85)
MUTED = RGBColor(100, 116, 139)
WHITE = RGBColor(248, 250, 252)
TEAL = RGBColor(20, 184, 166)
AMBER = RGBColor(245, 158, 11)
SKY = RGBColor(14, 165, 233)
SURFACE = RGBColor(241, 245, 249)
SURFACE_ALT = RGBColor(226, 232, 240)
SUCCESS = RGBColor(34, 197, 94)

TITLE_FONT = "Aptos Display"
BODY_FONT = "Aptos"
CODE_FONT = "Menlo"


def set_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_top_band(slide, accent=TEAL):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.28))
    band.fill.solid()
    band.fill.fore_color.rgb = accent
    band.line.fill.background()


def add_footer(slide, text="Bazel Techtalk"):  # noqa: D401
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(5.2), Inches(0.25))
    frame = box.text_frame
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED


def add_title(slide, title, subtitle=None, dark=False):
    color = WHITE if dark else NAVY
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(8.8), Inches(1.1))
    frame = title_box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = TITLE_FONT
    run.font.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = color
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.45), Inches(8.6), Inches(0.6))
        sub_frame = sub_box.text_frame
        sub_frame.word_wrap = True
        sub_p = sub_frame.paragraphs[0]
        sub_run = sub_p.add_run()
        sub_run.text = subtitle
        sub_run.font.name = BODY_FONT
        sub_run.font.size = Pt(15)
        sub_run.font.color.rgb = WHITE if dark else SLATE


def add_bullets(slide, bullets, left=0.9, top=1.9, width=6.7, height=3.9, color=NAVY):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = bullet
        run.font.name = BODY_FONT
        run.font.size = Pt(21)
        run.font.color.rgb = color


def add_emphasis_box(slide, text, left=8.0, top=2.0, width=4.3, height=2.3, fill=SLATE, text_color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.bold = True
    run.font.size = Pt(22)
    run.font.color.rgb = text_color


def add_code_box(slide, title, code_lines, left=0.9, top=2.0, width=11.5, height=3.9):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(left + 0.3), Inches(top + 0.15), Inches(width - 0.6), Inches(0.35))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_run = title_p.add_run()
    title_run.text = title
    title_run.font.name = BODY_FONT
    title_run.font.bold = True
    title_run.font.size = Pt(16)
    title_run.font.color.rgb = RGBColor(148, 163, 184)

    code_box = slide.shapes.add_textbox(Inches(left + 0.35), Inches(top + 0.55), Inches(width - 0.7), Inches(height - 0.75))
    code_frame = code_box.text_frame
    code_frame.word_wrap = True
    for index, line in enumerate(code_lines):
        p = code_frame.paragraphs[0] if index == 0 else code_frame.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = line
        run.font.name = CODE_FONT
        run.font.size = Pt(18)
        run.font.color.rgb = WHITE


def add_stat_card(slide, title, body, left, top, width=3.8, height=1.45, fill=SURFACE):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.fill.background()

    tf = card.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = BODY_FONT
    r1.font.bold = True
    r1.font.size = Pt(17)
    r1.font.color.rgb = NAVY

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = BODY_FONT
    r2.font.size = Pt(14)
    r2.font.color.rgb = SLATE


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_background(slide, NAVY)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(9.5), 0, Inches(3.833), Inches(7.5))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()
    accent2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(8.6), Inches(3.95), Inches(4.733), Inches(0.42))
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = AMBER
    accent2.line.fill.background()
    add_title(slide, "Bazel", "Explain the model first, then prove it live", dark=True)
    add_emphasis_box(slide, "Speed\nCorrectness\nScale", left=0.72, top=3.25, width=3.1, height=2.1, fill=RGBColor(30, 41, 59))
    add_footer(slide, "Techtalk deck generated from workspace content")

    content_slides = [
        {
            "title": "Bazel and Important Glossary",
            "bullets": [
                "Bazel: graph-based build and test system.",
                "Target: buildable or testable unit (for example //app:demo_app).",
                "Dependency graph: explicit relationships between targets.",
                "Cache: reuse outputs when inputs are unchanged.",
            ],
            "box": ("Focus", "Core terms first.", AMBER),
        },
        {
            "title": "Q&A Round 1",
            "bullets": [
                "Any questions on Bazel basics so far?",
                "Target, dependency graph, and cache are the key ideas.",
                "If this is clear, we will move to a quick concept check.",
            ],
            "box": ("Goal", "Quick clarity before demo.", TEAL),
        },
        {
            "title": "Quick Check Questions",
            "bullets": [
                "If one file changes, should everything rebuild?",
                "What does //app:demo_app represent?",
                "What does Bazel reuse when inputs are unchanged?",
                "Which command shows dependencies of a target?",
            ],
            "box": ("Engagement", "Check understanding.", SKY),
        },
        {
            "title": "Simple Project Demo",
            "bullets": [
                "Use app/, lib/, and tests/ to show target structure.",
                "Run query, run, test, and small edit workflow.",
                "Call out selective rebuild and cache behavior.",
            ],
            "box": ("Approach", "Show it live.", SUCCESS),
        },
        {
            "title": "Q&A Round 2",
            "bullets": [
                "Any questions from the live demo so far?",
                "We can revisit deps and rdeps if you want impact examples.",
                "If this is clear, we will move to hands-on practice.",
            ],
            "box": ("Goal", "Answer practical questions.", AMBER),
        },
        {
            "title": "Audience Hands-On",
            "bullets": [
                "Clone https://github.com/Ganagood/Techtalk.git to your machine.",
                "Open the repo locally and run the same Bazel commands.",
                "Repeat the flow and see the rebuild behavior yourself.",
            ],
            "box": ("Practice", "Audience runs it now.", SUCCESS),
        },
        {
            "title": "Official Documentation Guide",
            "bullets": [
                "Main docs: https://bazel.build",
                "Query reference: https://bazel.build/query/language",
                "Bzlmod guide: https://bazel.build/external/overview",
                "Rules and ecosystem: https://github.com/bazelbuild",
            ],
            "box": ("Next", "Share official docs.", SKY),
        },
        {
            "title": "Fun Close",
            "bullets": [
                "Funny question: Is your build faster than your coffee machine?",
                "Funny question: How many times per day do we rebuild unchanged code?",
                "If no doubts remain, close the session with thanks.",
            ],
            "box": ("Wrap-up", "Light close.", TEAL),
        },
    ]

    for item in content_slides:
        slide = prs.slides.add_slide(blank)
        set_background(slide)
        add_top_band(slide)
        add_title(slide, item["title"])
        add_bullets(slide, item["bullets"])
        label, body, color = item["box"]
        add_emphasis_box(slide, f"{label}\n{body}", fill=color)
        add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_background(slide, RGBColor(248, 250, 252))
    add_top_band(slide, AMBER)
    add_title(slide, "Live Demo Commands", "Watch how Bazel responds to each command")
    add_code_box(
        slide,
        "Core demo flow",
        [
            "bazel query //...",
            "bazel run //app:demo_app",
            "bazel run //app:demo_app",
            "bazel test //tests:math_utils_test",
            "# edit lib/math_utils.cc: 1.07 -> 1.08",
            "bazel run //app:demo_app",
            "bazel query 'deps(//app:demo_app)'",
        ],
        top=1.85,
        height=4.45,
    )
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_background(slide, RGBColor(248, 250, 252))
    add_top_band(slide, TEAL)
    add_title(slide, "Audience Practice", "Clone the repo and run the same commands locally")
    add_code_box(
        slide,
        "Local practice flow",
        [
            "git clone https://github.com/Ganagood/Techtalk.git",
            "cd Techtalk/demo-project",
            "bazel query //...",
            "bazel run //app:demo_app",
            "bazel test //tests:math_utils_test",
            "# edit lib/math_utils.cc: 0.07 -> 0.08",
            "bazel run //app:demo_app",
        ],
        top=1.85,
        height=4.1,
    )
    add_stat_card(slide, "Repo", "https://github.com/Ganagood/Techtalk.git", 0.95, 6.1, width=11.4, height=0.95, fill=SURFACE_ALT)
    add_footer(slide)

    slide = prs.slides.add_slide(blank)
    set_background(slide, NAVY)
    add_top_band(slide, SUCCESS)
    add_title(slide, "Close", "Any doubts before we wrap?", dark=True)
    add_stat_card(slide, "Funny check", "If Bazel saves 15 minutes per engineer per day, what is that worth for your whole team in a month?", 0.9, 2.2, width=11.4, height=1.2, fill=RGBColor(30, 41, 59))
    add_stat_card(slide, "No more questions", "If everything is clear, thank the audience and invite them to continue in the official docs.", 0.9, 3.75, width=11.4, height=1.15, fill=RGBColor(30, 41, 59))
    add_stat_card(slide, "Final line", "Bazel helps teams spend less time waiting and more time building.", 0.9, 5.2, width=11.4, height=1.0, fill=RGBColor(30, 41, 59))
    add_footer(slide, "Thank you | Questions are welcome")

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build_deck()
    print(path)
